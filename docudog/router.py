"""Routing: filters, hashing, text extraction, dedupe, inference + reporting."""

from __future__ import annotations

import hashlib
import logging
import os
from collections.abc import Callable
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from . import activity
from . import audit
from . import categories
from . import context_bundles
from . import extract_hwp
from . import inference
from . import last_classify
from . import owner_tags
from . import paths_util
from . import power_gate
from . import related_docs
from . import reporter
from . import rule_hints
from . import semantic_diff
from . import skip_insights
from . import status_dashboard

logger = logging.getLogger(__name__)

_SKIPPABLE_FOR_MVP = {".pdf"}


def _normalize_path(path: str) -> str:
    return paths_util.normalize_fs_path(path)


def passes_file_filters(config: dict[str, Any], path: str) -> bool:
    filters = config.get("file_filters", {})
    exts = {e.lower() for e in filters.get("allowed_extensions", [])}
    ext = Path(path).suffix.lower()
    if ext not in exts:
        return False
    size = filters.get("size_limit", {})
    min_b = int(size.get("min_bytes", 0))
    max_b = int(size.get("max_bytes", 2**62))
    try:
        sz = os.path.getsize(path)
    except OSError:
        return False
    return min_b <= sz <= max_b


def sha256_file(
    path: str,
    chunk_size: int = 1024 * 1024,
    cfg: dict[str, Any] | None = None,
) -> str:
    def _once() -> str:
        h = hashlib.sha256()
        with open(path, "rb") as f:
            while True:
                chunk = f.read(chunk_size)
                if not chunk:
                    break
                h.update(chunk)
        return h.hexdigest()

    return paths_util.with_file_retry(cfg, f"sha256:{path}", _once)


def _read_plain_text(path: str) -> str:
    with open(path, encoding="utf-8", errors="replace") as f:
        return f.read()


def _read_docx(path: str) -> str:
    import docx  # type: ignore[import-untyped]

    document = docx.Document(path)
    parts: list[str] = []
    for p in document.paragraphs:
        t = (p.text or "").strip()
        if t:
            parts.append(t)
    for table in document.tables:
        for row in table.rows:
            cells = " | ".join((c.text or "").strip() for c in row.cells)
            if cells.strip():
                parts.append(cells)
    return "\n".join(parts)


def _read_pptx(path: str) -> str:
    from pptx import Presentation  # type: ignore[import-untyped]

    prs = Presentation(path)
    out: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            t = (shape.text or "").strip()
            if t:
                out.append(t)
    return "\n".join(out)


def _read_xlsx(path: str) -> str:
    from openpyxl import load_workbook  # type: ignore[import-untyped]

    wb = load_workbook(filename=path, read_only=True, data_only=True)
    lines: list[str] = []
    try:
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None and str(c).strip()]
                if cells:
                    lines.append(" | ".join(cells))
    finally:
        wb.close()
    return "\n".join(lines)


def _read_hwp(path: str) -> str:
    return extract_hwp.read_hwp_text(path)


def extract_document_text(
    path: str,
    cfg: dict[str, Any] | None = None,
) -> tuple[str | None, str | None]:
    """
    Return (text, skip_reason). skip_reason is set when routing should stop
    without treating as an error (e.g., MVP skips PDF).
    """
    ext = Path(path).suffix.lower()
    if ext in _SKIPPABLE_FOR_MVP:
        return None, f"MVP 스킵: 이 확장자는 텍스트 추출 미구현 ({ext})"

    def _once() -> tuple[str | None, str | None]:
        if ext in {".txt", ".md"}:
            return _read_plain_text(path), None
        if ext == ".docx":
            return _read_docx(path), None
        if ext == ".pptx":
            return _read_pptx(path), None
        if ext == ".xlsx":
            return _read_xlsx(path), None
        if ext in {".hwp", ".hwpx"}:
            try:
                return _read_hwp(path), None
            except extract_hwp.EncryptedHwpError as e:
                return None, str(e)
        return None, f"지원되지 않는 확장자: {ext}"

    try:
        return paths_util.with_file_retry(cfg, f"extract:{path}", _once)
    except Exception as e:
        logger.debug("텍스트 추출 실패: %s — %s", path, e)
        logger.debug("Extract traceback", exc_info=True)
        return None, f"텍스트 추출 실패: {e}"


def _utc_now_iso() -> str:
    return datetime.now(timezone.utc).isoformat()


def process_file(
    config: dict[str, Any],
    state: dict[str, Any],
    path: str,
    report_path: str,
    state_path: str,
    idle_trigger_seconds: float,
    seconds_since_last_input: Callable[[], float],
    save_state: Callable[[], None],
    file_event_unix: float | None = None,
    get_fs_event_snapshot: Callable[[], list[tuple[str, float]]] | None = None,
) -> str | None:
    """
    Full routing for one file: filters, extract, hash, dedupe, inference, report.

    state_path: used to locate DocuDog_tag_overrides.json next to state.

    Returns:
        None — finished handling (including filters/skips and successful analysis).
        "requeue" — same path should be enqueued again (user active / yielded).
        "requeue_power" — power/battery gate; requeue and pause drain briefly.

    Yield / pause: if the user becomes active during streaming inference, we stop
    without writing partial results; caller should requeue the path.
    """
    norm = _normalize_path(path)
    files_state: dict[str, Any] = state.setdefault("files", {})

    if not passes_file_filters(config, norm):
        logger.debug("Skip (filter): %s", norm)
        activity.append_activity(config, report_path, "skip_filter", norm)
        return None

    text, skip_reason = extract_document_text(norm, config)
    if skip_reason:
        logger.debug("Skip: %s — %s", norm, skip_reason)
        insight = skip_insights.record_extract_skip(state, config, norm, skip_reason)
        save_state()
        note = f"{skip_reason} — `{norm}`"
        if insight.get("sensitive_keyword"):
            note = (
                f"민감 파일명 키워드({insight['sensitive_keyword']}) · {note}"
            )
        reporter.append_note(report_path, note)
        reporter.refresh_report_banner(report_path, state)
        activity.append_activity(
            config,
            report_path,
            "skip_extract",
            f"{norm} | {skip_reason}"
            + (
                f" | sensitive={insight['sensitive_keyword']}"
                if insight.get("sensitive_keyword")
                else ""
            ),
        )
        try:
            status_dashboard.write_status(config, state, report_path, save_state)
        except Exception:
            logger.exception("Status dashboard update failed after skip")
        return None
    if not (text and text.strip()):
        logger.debug("Skip (empty text): %s", norm)
        activity.append_activity(config, report_path, "skip_empty", norm)
        return None

    digest = sha256_file(norm, cfg=config)
    prev = files_state.get(norm)

    if isinstance(prev, dict) and prev.get("sha256") == digest:
        prev["last_checked_utc"] = _utc_now_iso()
        files_state[norm] = prev
        save_state()
        logger.debug("Skip LLM (unchanged content hash): %s", norm)
        activity.append_activity(config, report_path, "skip_hash", norm)
        return None

    def user_active() -> bool:
        return seconds_since_last_input() < float(idle_trigger_seconds)

    if user_active():
        logger.info("Defer (user active before inference): %s", norm)
        activity.append_activity(config, report_path, "defer_active", norm)
        return "requeue"

    rule_hits = rule_hints.match_rules(text, config)
    hint_parts: list[str] = []
    rb = rule_hints.prompt_hint_block(rule_hits)
    if rb:
        hint_parts.append(rb)
    cats_doc = categories.load_categories(config, state_path)
    cb = categories.prompt_block(cats_doc, config)
    if cb:
        hint_parts.append(cb)
    rule_block = "\n\n".join(hint_parts)

    power_ok, power_reason = power_gate.inference_power_allowed(config)
    if not power_ok:
        logger.info("Defer (power gate): %s — %s", norm, power_reason)
        activity.append_activity(
            config,
            report_path,
            "defer_power",
            f"{norm} | {power_reason}",
        )
        return "requeue_power"

    try:
        result = inference.classify_document(
            text,
            config,
            should_yield=user_active,
            rule_hint=rule_block,
        )
    except inference.YieldToUser:
        logger.info("Yielded during inference (user active): %s", norm)
        activity.append_activity(config, report_path, "defer_yield", norm)
        return "requeue"

    inf_src = str(result.pop(inference.DOCUDOG_META_SOURCE, "mock"))
    inf_reason = str(result.pop(inference.DOCUDOG_META_REASON, ""))

    tags = result.get("tags") or []
    if isinstance(tags, str):
        tags = [tags]
    sec = str(result.get("security_level", ""))
    summary = str(result.get("summary", ""))

    model_tags = list(tags)
    model_sec = sec
    overrides_doc = owner_tags.load_tag_overrides(config, state_path)
    tags, sec, owner_applied = owner_tags.merge_owner_tags(
        norm, model_tags, model_sec, overrides_doc
    )
    floor_applied: str | None = None
    if not owner_applied:
        sec, floor_applied = rule_hints.apply_security_floor(sec, rule_hits)
        if floor_applied:
            note = f"rule_floor={floor_applied}"
            inf_reason = f"{inf_reason}; {note}".strip("; ") if inf_reason else note

    category_ids, cat_needs_review = categories.resolve_from_result(
        result, cats_doc, config
    )
    # strip category keys from result leftovers (already used)
    result.pop("category_id", None)
    result.pop("category_ids", None)
    report_tags = list(tags)
    if category_ids:
        labels = [
            categories.label_for(c, cats_doc) for c in category_ids if c != "uncategorized"
        ]
        if labels:
            report_tags = [f"[{', '.join(labels)}]"] + list(tags)
        elif "uncategorized" in category_ids:
            report_tags = ["[uncategorized]"] + list(tags)
        if cat_needs_review:
            note = "category_needs_review"
            inf_reason = f"{inf_reason}; {note}".strip("; ") if inf_reason else note

    change_line = ""
    if (
        semantic_diff.semantic_enabled(config)
        and isinstance(prev, dict)
        and prev.get("sha256")
        and prev.get("sha256") != digest
    ):
        try:
            change_line = semantic_diff.build_change_line(
                str(prev.get("summary") or ""),
                summary,
                "",  # previous full text not retained; summary-driven
                text[:2000],
                config,
            )
        except Exception:
            logger.debug("semantic change line failed", exc_info=True)

    analyzed_at = datetime.now(timezone.utc)
    reporter.append_classification(
        report_path=report_path,
        analyzed_at_utc=analyzed_at,
        file_name=os.path.basename(norm),
        file_hash=digest,
        tags=report_tags,
        security_level=sec,
        summary=summary,
        inference_source=inf_src,
        inference_reason=inf_reason,
        owner_tags_applied=owner_applied,
        config=config,
        state=state,
    )

    excerpt_lim = 1200
    acfg = config.get("audit_settings")
    if isinstance(acfg, dict) and acfg.get("max_excerpt_chars") is not None:
        try:
            excerpt_lim = int(acfg["max_excerpt_chars"])
        except (TypeError, ValueError):
            excerpt_lim = 1200
    audit.maybe_record_sensitive_classification(
        config,
        norm_path=norm,
        report_path=report_path,
        analyzed_at_utc=analyzed_at,
        file_hash=digest,
        tags=tags,
        security_level=sec,
        text_excerpt=text[: max(0, excerpt_lim)],
        inference_source=inf_src,
    )

    files_state[norm] = {
        "sha256": digest,
        "last_analyzed_utc": analyzed_at.isoformat(),
        "last_checked_utc": _utc_now_iso(),
        "summary": summary,
        "model_tags": model_tags,
        "model_security_level": model_sec,
        "tags": tags,
        "security_level": sec,
        "owner_override": owner_applied,
        "inference_source": inf_src,
        "inference_reason": inf_reason,
        "rule_hits": [h.get("id") for h in rule_hits],
        "rule_floor": floor_applied or "",
        "category_ids": category_ids,
        "category_needs_review": cat_needs_review,
    }
    if isinstance(prev, dict) and isinstance(prev.get("summary_history"), list):
        files_state[norm]["summary_history"] = list(prev["summary_history"])
    if semantic_diff.semantic_enabled(config):
        sset = config.get("semantic_settings")
        max_h = 8
        if isinstance(sset, dict) and sset.get("history_max") is not None:
            try:
                max_h = int(sset["history_max"])
            except (TypeError, ValueError):
                max_h = 8
        semantic_diff.push_summary_history(
            files_state[norm],
            sha256=digest,
            summary=summary,
            utc=analyzed_at.isoformat(),
            change_summary=change_line,
            max_entries=max_h,
        )
    # related candidates (same key / summary / context bundles)
    related = related_docs.find_related_paths(state, norm, files_state[norm])
    if related:
        files_state[norm]["related_paths"] = related
        state["last_related"] = {
            "anchor": norm,
            "paths": related,
            "utc": analyzed_at.isoformat(),
        }
    state["last_inference_backend"] = inf_src
    state["last_inference_utc"] = analyzed_at.isoformat()
    save_state()
    act_msg = f"{norm} | {sec} | {inf_src}"
    if category_ids:
        act_msg += f" | cat={','.join(category_ids)}"
    if change_line:
        act_msg += f" | change={change_line[:80]}"
    activity.append_activity(
        config,
        report_path,
        "classify",
        act_msg,
        when=analyzed_at,
    )
    try:
        last_classify.write_last_classify(
            config,
            report_path,
            path=norm,
            security_level=sec,
            tags=list(tags) if isinstance(tags, list) else [str(tags)],
            summary=summary,
            inference_source=inf_src,
            file_hash=digest,
            when=analyzed_at,
            category_ids=category_ids,
            change_summary=change_line,
        )
    except Exception:
        logger.exception("last_classify write failed for %s", norm)
    logger.info(
        "Analyzed: %s [inference=%s]",
        os.path.basename(norm),
        inf_src,
    )
    if get_fs_event_snapshot is not None and file_event_unix is not None:
        try:
            snap = get_fs_event_snapshot()
            context_bundles.update_bundles_after_analysis(
                config, state, norm, float(file_event_unix), snap, save_state
            )
        except Exception:
            logger.exception("Context bundle update failed for %s", norm)
    try:
        status_dashboard.write_status(config, state, report_path, save_state)
    except Exception:
        logger.exception("Status dashboard update failed for %s", norm)
    return None
